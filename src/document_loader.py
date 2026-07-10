"""文档加载器 — 从知识库目录读取各类文件，提取纯文本

支持格式: PDF, DOCX, PPTX, XLSX, TXT, PNG, JPG, JPEG
扫描件 PDF 自动调用 PP-OCRv4 进行 OCR 识别
"""
import os
import re
import json
import subprocess
from pathlib import Path
from typing import List, Dict, Optional

from src.config_loader import config


# ============================================================
# OCR 辅助 — 调用 PP-OCR 子进程
# ============================================================

_OCR_SCRIPT = Path(__file__).parent.parent / "scripts" / "ocr_worker.py"
_OCR_PYTHON = "python3"  # 全局 Python，已安装 PaddleOCR


def _ocr_image(image_path: str) -> Optional[str]:
    """调用 PP-OCR 子进程识别图片文本"""
    try:
        result = subprocess.run(
            [_OCR_PYTHON, str(_OCR_SCRIPT), image_path],
            capture_output=True, text=True, timeout=120,
            stderr=subprocess.DEVNULL
        )
        # 从 stdout 中找到最后一行 JSON（跳过 PaddleOCR 的日志行）
        lines = result.stdout.strip().split('\n')
        json_line = ''
        for line in reversed(lines):
            line = line.strip()
            if line.startswith('{'):
                json_line = line
                break
        if not json_line:
            print(f"  [WARN] OCR 输出中没有 JSON")
            return None
        data = json.loads(json_line)
        if data.get("success"):
            return data["text"]
        else:
            print(f"  [WARN] OCR 失败: {data.get('error', '未知错误')}")
            return None
    except subprocess.TimeoutExpired:
        print(f"  [WARN] OCR 超时")
        return None
    except Exception as e:
        print(f"  [WARN] OCR 调用异常: {e}")
        return None


# ============================================================
# 单文件解析
# ============================================================

def _clean_pdf_text(text: str) -> str:
    """清理 PDF 文本中的页码等噪声"""
    # 去掉行首的 ". 117" 页码（可能接正文）
    text = re.sub(r'^\.\s*\d{1,4}\s*', '', text, flags=re.MULTILINE)
    # 去掉行首的纯数字页码
    text = re.sub(r'^\d{1,4}\s*$', '', text, flags=re.MULTILINE)
    # 去掉 "- 117 -" 格式
    text = re.sub(r'^-\s*\d{1,4}\s*-\s*$', '', text, flags=re.MULTILINE)
    # 清理多余空行
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def extract_text_from_pdf(filepath: str) -> Optional[str]:
    """提取 PDF 文本（PyMuPDF + 扫描件自动 OCR）"""
    try:
        import fitz
        doc = fitz.open(filepath)
        text_parts = []
        for page_num, page in enumerate(doc):
            t = page.get_text()
            text_len = len(t.strip()) if t else 0

            # 如果文本很少（< 50 字），视为扫描件/图片页，调用 OCR
            if text_len < 50:
                pix = page.get_pixmap(dpi=200)
                img_path = f"/tmp/ocr_page_{Path(filepath).stem}_{page_num}.png"
                pix.save(img_path)
                ocr_text = _ocr_image(img_path)
                try:
                    os.remove(img_path)
                except OSError:
                    pass
                if ocr_text and ocr_text.strip():
                    text_parts.append(ocr_text.strip())
                continue

            # 正常提取的文本
            cleaned = _clean_pdf_text(t.strip())
            if cleaned:
                text_parts.append(cleaned)

        doc.close()
        if text_parts:
            return "\n\n".join(text_parts)
    except Exception as e:
        print(f"  [WARN] PyMuPDF 提取失败 ({Path(filepath).name}): {e}")

    # 方案2: pypdf 回退
    try:
        import pypdf
        reader = pypdf.PdfReader(filepath)
        text_parts = []
        for page in reader.pages:
            t = page.extract_text()
            if t and t.strip():
                text_parts.append(t.strip())
        return "\n".join(text_parts) if text_parts else None
    except Exception as e:
        print(f"  [WARN] PDF 提取失败 ({Path(filepath).name}): {e}")
        return None


def extract_text_from_image(filepath: str) -> Optional[str]:
    """提取图片文本（调用 PP-OCR）"""
    return _ocr_image(filepath)


def extract_text_from_docx(filepath: str) -> Optional[str]:
    """提取 DOCX 文本"""
    try:
        from docx import Document
        doc = Document(filepath)
        texts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n".join(texts) if texts else None
    except Exception as e:
        print(f"  [WARN] DOCX 提取失败 ({Path(filepath).name}): {e}")
        return None


def extract_text_from_doc(filepath: str) -> Optional[str]:
    """提取旧版 DOC 文本（尝试用 python-docx, 否则返回提示）"""
    try:
        from docx import Document
        doc = Document(filepath)
        texts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n".join(texts) if texts else None
    except Exception:
        print(f"  [INFO] 旧版 .doc 文件，建议用 Word 另存为 .docx: {Path(filepath).name}")
        return None


def extract_text_from_pptx(filepath: str) -> Optional[str]:
    """提取 PPTX 文本（所有幻灯片的文本框）"""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
        return "\n".join(texts) if texts else None
    except Exception as e:
        print(f"  [WARN] PPTX 提取失败 ({Path(filepath).name}): {e}")
        return None


def extract_text_from_xlsx(filepath: str) -> Optional[str]:
    """提取 XLSX 文本（合并所有sheet的单元格文本）"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
        texts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_texts = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if row_texts:
                    texts.append(" | ".join(row_texts))
        wb.close()
        return "\n".join(texts) if texts else None
    except Exception as e:
        print(f"  [WARN] XLSX 提取失败 ({Path(filepath).name}): {e}")
        return None


def extract_text_from_txt(filepath: str) -> Optional[str]:
    """提取 TXT 文本（自动检测编码）"""
    encodings = ["utf-8", "gbk", "gb2312", "gb18030", "utf-16"]
    for enc in encodings:
        try:
            with open(filepath, "r", encoding=enc) as f:
                text = f.read().strip()
            return text if text else None
        except (UnicodeDecodeError, UnicodeError):
            continue
    print(f"  [WARN] 无法解码 TXT: {Path(filepath).name}")
    return None


# ============================================================
# 文件分派器
# ============================================================

EXTENSION_MAP = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".doc": extract_text_from_doc,
    ".pptx": extract_text_from_pptx,
    ".xlsx": extract_text_from_xlsx,
    ".xls": extract_text_from_xlsx,
    ".txt": extract_text_from_txt,
    ".png": extract_text_from_image,
    ".jpg": extract_text_from_image,
    ".jpeg": extract_text_from_image,
}


def extract_text(filepath: str) -> Optional[str]:
    """根据扩展名自动调度提取器"""
    ext = Path(filepath).suffix.lower()
    extractor = EXTENSION_MAP.get(ext)
    if extractor is None:
        print(f"  [SKIP] 不支持的文件格式: {filepath}")
        return None
    return extractor(filepath)


# ============================================================
# 批量扫描知识库
# ============================================================

def scan_documents(kb_path: Optional[str] = None) -> List[Dict]:
    """
    扫描知识库目录，返回文档元信息列表。

    返回:
    [
        {
            "filepath": "绝对路径",
            "relpath": "相对知识库根目录的路径",
            "filename": "文件名",
            "category": "顶层分类",
            "ext": ".pdf"
        },
        ...
    ]
    """
    if kb_path is None:
        kb_path = config["documents"]["knowledge_base"]

    kb_root = Path(kb_path)
    if not kb_root.exists():
        raise FileNotFoundError(f"知识库目录不存在: {kb_root}")

    supported = config["documents"]["supported_extensions"]
    # 确保图片格式也被支持
    for img_ext in [".png", ".jpg", ".jpeg", ".bmp"]:
        if img_ext not in supported:
            supported.append(img_ext)

    docs = []

    for item in kb_root.rglob("*"):
        if not item.is_file():
            continue
        ext = item.suffix.lower()
        if ext not in supported:
            continue

        filename = item.name

        # 排除临时文件和系统文件
        if filename.startswith("~$") or filename.startswith("."):
            continue
        if filename.lower() in ("thumbs.db", "desktop.ini"):
            continue

        # 计算相对路径（只取第一层子目录作为 category）
        rel = item.relative_to(kb_root)
        parts = rel.parts
        category = parts[0] if len(parts) > 1 else "未分类"

        docs.append({
            "filepath": str(item.resolve()),
            "relpath": str(rel),
            "filename": filename,
            "category": category,
            "ext": ext,
        })

    docs.sort(key=lambda d: d["relpath"])
    return docs


def load_all_documents(kb_path: Optional[str] = None, verbose: bool = True) -> List[Dict]:
    """
    扫描并提取所有文档的文本内容。

    返回:
    [
        {
            "filepath": "...",
            "relpath": "...",
            "filename": "...",
            "category": "...",
            "text": "提取的纯文本",
            "char_count": 1234,
        },
        ...
    ]
    """
    docs = scan_documents(kb_path)

    results = []
    for i, doc in enumerate(docs):
        if verbose:
            print(f"[{i+1}/{len(docs)}] 提取: {doc['relpath']}")
        text = extract_text(doc["filepath"])
        if text and len(text.strip()) > 10:  # 至少 10 个有效字符
            doc["text"] = text
            doc["char_count"] = len(text)
            results.append(doc)
        else:
            if verbose:
                print(f"  -> 跳过（内容为空或太少）")

    if verbose:
        print(f"\n成功加载 {len(results)}/{len(docs)} 个文档")

    return results


# ============================================================
# 测试入口
# ============================================================
if __name__ == "__main__":
    docs = scan_documents()
    print(f"发现 {len(docs)} 个文档:")
    for d in docs[:10]:
        print(f"  [{d['category']}] {d['relpath']}")
    print("...")
    print("\n提取全部文本...")
    loaded = load_all_documents(verbose=True)
    print(f"\n总字符数: {sum(d['char_count'] for d in loaded)}")